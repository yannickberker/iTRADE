"""Handle export to PowerPoint files."""
import io
from os import PathLike
from typing import TypedDict

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE


# https://github.com/scanny/python-pptx/issues/333
# https://github.com/scanny/python-pptx/pull/439
def monkey_patch_pptx_picture_in_placeholder() -> None:
    """Monkey patch python-pptx for a better insert_pictures in (Slide)Placeholders."""
    # pylint: disable=import-outside-toplevel, protected-access
    from pptx.oxml.shapes.picture import CT_Picture
    from pptx.shapes.shapetree import (
        PicturePlaceholder,
        PlaceholderPicture,
        SlidePlaceholder,
    )

    def insert_picture(
        self: PlaceholderPicture,
        image_file: str | io.BytesIO,
        crop: bool = False,
        vcenter: bool = True,
        hcenter: bool = True,
    ) -> PlaceholderPicture:
        """See PicturePlaceholder.insert_picture."""
        # https://github.com/scanny/python-pptx/blob/v0.6.21/pptx
        # .../shapes/placeholder.py#L310-L321
        pic = self._new_placeholder_pic(image_file, crop, vcenter, hcenter)
        self._replace_placeholder_with(pic)
        return PlaceholderPicture(pic, self._parent)

    def _new_placeholder_pic(
        self: PlaceholderPicture,
        image_file: str | io.BytesIO,
        crop: bool = False,
        vcenter: bool = True,
        hcenter: bool = True,
    ) -> CT_Picture:
        """See PicturePlaceholder._new_placeholder_pic."""
        # https://github.com/scanny/python-pptx/blob/v0.6.21/pptx
        # .../shapes/placeholder.py#L323-L334
        r_id, desc, image_size = self._get_or_add_image(image_file)
        shape_id, name = self.shape_id, self.name
        if crop:
            pic = CT_Picture.new_ph_pic(shape_id, name, desc, r_id)
            pic.crop_to_fit(image_size, (self.width, self.height))
        else:
            ph_w, ph_h = self.width, self.height
            aspect_ph = ph_w / ph_h

            img_w, img_h = image_size
            aspect_img = img_w / img_h

            left = self.left
            top = self.top

            if aspect_ph > aspect_img:
                w = int(ph_h * aspect_img)
                h = ph_h

                if hcenter:
                    left = left + (ph_w - w) / 2
            else:
                w = ph_w
                h = int(ph_w / aspect_img)

                if vcenter:
                    top = top + (ph_h - h) / 2

            pic = CT_Picture.new_pic(shape_id, name, desc, r_id, left, top, w, h)
        return pic

    SlidePlaceholder.insert_picture = insert_picture
    SlidePlaceholder._new_placeholder_pic = _new_placeholder_pic
    SlidePlaceholder._get_or_add_image = PicturePlaceholder._get_or_add_image


# https://github.com/scanny/python-pptx/issues/702
def monkey_patch_pptx_no_last_modified_date() -> None:
    """Monkey patch python-pptx to write zip without last-modified dates."""
    # pylint: disable=import-outside-toplevel, protected-access
    from zipfile import ZipInfo

    from pptx.opc.serialized import _ZipPkgWriter

    _ZipPkgWriter.write = lambda self, pack_uri, blob: self._zipf.writestr(
        ZipInfo(pack_uri.membername),
        blob,
        compress_type=self._zipf.compression,
        compresslevel=self._zipf.compresslevel,
    )


class TextAndImage(TypedDict):
    """Content for one PowerPoint slide."""

    text: str
    image: str | io.BytesIO


def write_powerpoint(
    contents: list[TextAndImage],
    filename: PathLike,
    pptx: PathLike | None = None,
) -> None:
    """Save a list of slide content dictionaries into a PowerPoint presentation."""
    monkey_patch_pptx_picture_in_placeholder()
    monkey_patch_pptx_no_last_modified_date()

    prs = Presentation(pptx=pptx)

    for content in contents:
        text = content.get("text")
        image = content.get("image")

        layout = prs.slide_layouts.get_by_name("Title and Content")
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = text
        slide.shapes.title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        slide.placeholders[1].insert_picture(image)

    prs.save(filename)
